import os
import mimetypes
import shutil
from flask import Flask, render_template, request, redirect, url_for, flash, send_file
from werkzeug.utils import secure_filename
import mammoth
from pptx import Presentation
import fitz  # PyMuPDF
import nbformat
from jupyter_client import KernelManager

app = Flask(__name__, static_url_path='/static', static_folder='static')
app.config['UPLOAD_FOLDER'] = os.path.abspath('templates/uploads')
app.secret_key = 'supersecretkey'
file_index = {}

# Allowed file extensions
app.config['ALLOWED_EXTENSIONS'] = {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'docx', 'pptx', 'ipynb'}


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def create_directory(directory):
    # Create the directory if it doesn't exist
    if not os.path.exists(directory):
        os.makedirs(directory)
        print(f"Created directory: {directory}")


def index_files(upload_folder):
    file_index.clear()
    for root, dirs, files in os.walk(upload_folder):
        for file in files:
            file_path = os.path.join(root, file)
            content = None
            if file.endswith('.txt'):
                content = extract_text_from_txt(file_path)
            elif file.endswith('.pdf'):
                content = extract_text_from_pdf(file_path)
            elif file.endswith('.docx'):
                content = extract_text_from_docx(file_path)
            elif file.endswith('.pptx'):
                content = extract_text_from_pptx(file_path)
            elif file.endswith('.ipynb'):
                content = extract_text_from_ipynb(file_path)
            if content is not None:
                relative_path = os.path.relpath(file_path, upload_folder)
                file_index[relative_path] = content
                print(f'Indexed file: {relative_path}')
                new_relative_path = relative_path.replace(" ", "")
                if new_relative_path != relative_path:
                    file_index[new_relative_path] = content
                    print(f'Reindexed file: {new_relative_path}')


def extract_text_from_txt(path):
    with open(path, 'r', encoding='utf-8') as file:
        return file.read()


def extract_text_from_pdf(path):
    doc = fitz.open(path)
    text = []
    for page in doc:
        text.append(page.get_text())
    return "\n".join(text)


def extract_text_from_docx(path):
    with open(path, 'rb') as file:
        result = mammoth.extract_raw_text(file)
        return result.value


def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_text_from_ipynb(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        nb = nbformat.read(f, as_version=4)
        text = ''
        for cell in nb['cells']:
            if cell['cell_type'] == 'markdown' or cell['cell_type'] == 'code':
                text += ''.join(cell['source']) + '\n\n'
        return text


def search_files(query):
    results = {}
    for file_path, content in file_index.items():
        if query.lower() in content.lower():
            results[file_path] = content
    return results


@app.route('/')
def index():
    categories = [d for root, dirs, files in os.walk(app.config['UPLOAD_FOLDER']) for d in dirs]
    return render_template('index.html', categories=categories)


@app.route('/pptx_player/<path:file_path>')
def pptx_player(file_path):
    absolute_file_path = os.path.abspath(os.path.join(app.config['UPLOAD_FOLDER'], file_path.replace('\\', '/')))
    if os.path.exists(absolute_file_path):
        return render_template('pptx_player.html', pptx_file=file_path)
    else:
        return "File not found", 404


@app.route('/search', methods=['POST'])
def search():
    query = request.form['query']
    search_results = search_files(query)
    return render_template('search_results.html', query=query, results=search_results)


@app.route('/pptx_preview')
def pptx_preview():
    slides_content = get_slides_content()
    return render_template('pptx_preview.html', slides_content=slides_content)


@app.route('/open/<path:file_path>', methods=['GET'])
def open_file(file_path):
    file_path = file_path.replace("\\", "/")
    absolute_file_path = os.path.abspath(os.path.join(app.config['UPLOAD_FOLDER'], file_path.replace('\\', '/')))
    print(f"Opening file: {absolute_file_path}")

    if os.path.exists(absolute_file_path):
        file_extension = file_path.rsplit('.', 1)[1].lower()
        if file_extension == 'pptx':
            images_dir = os.path.join('static', 'images')
            slides_content = extract_text_and_images_from_pptx(absolute_file_path, images_dir)
            slides_count = len(slides_content)
            print(f"Slides content: {slides_content}")
            return render_template('pptx_preview.html', slides_content=slides_content,
                                   slides_count=slides_count)
        elif file_extension == 'docx':
            with open(absolute_file_path, 'rb') as f:
                content = mammoth.extract_raw_text(f).value
            return render_template('preview_docx.html', content=content)
        elif file_extension == 'ipynb':
            with open(absolute_file_path, 'r', encoding='utf-8') as f:
                notebook = nbformat.read(f, as_version=4)
            return render_template('preview_ipynb.html', notebook=notebook)
        else:
            mime_type, _ = mimetypes.guess_type(absolute_file_path)
            if mime_type is not None:
                if mime_type.startswith('text'):
                    with open(absolute_file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    return render_template('preview_text.html', content=content)
                elif mime_type == 'application/pdf':
                    return send_file(absolute_file_path)
                elif mime_type.startswith('image'):
                    return render_template('preview_image.html', file_path=file_path)
                elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    return render_template('preview_docx.html', file_path=absolute_file_path)
                else:
                    return "File type not supported for preview"
            else:
                return "Unknown file type"
    else:
        return "File not found", 404


@app.route('/preview_ipynb/<path:file_path>', methods=['GET'])
def preview_ipynb(file_path):
    absolute_file_path = os.path.abspath(os.path.join(app.config['UPLOAD_FOLDER'], file_path))
    if os.path.exists(absolute_file_path):
        with open(absolute_file_path, 'r', encoding='utf-8') as f:
            notebook = nbformat.read(f, as_version=4)
        print(notebook)  # Debug: Print the notebook content
        return render_template('preview_ipynb.html', notebook=notebook)
    else:
        return "File not found", 404


@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        flash('No file part')
        return redirect(request.url)

    file = request.files['file']

    if file.filename == '':
        flash('No selected file')
        return redirect(request.url)

    if file and allowed_file(file.filename):
        category = request.form['category']
        category_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(category))

        if not os.path.exists(category_path):
            os.makedirs(category_path)

        file_path = os.path.join(category_path, secure_filename(file.filename))
        file.save(file_path)

        if file.filename.endswith('.pptx'):
            images_dir = os.path.join('static', 'images')
            slides_content = extract_text_and_images_from_pptx(file_path, images_dir)
            save_slides_by_category(slides_content, category_path)

        index_files(app.config['UPLOAD_FOLDER'])
        flash('File successfully uploaded')
        return redirect(url_for('index'))
    else:
        flash('File type not allowed')
        return redirect(request.url)


def save_slides_by_category(slides_content, category_path):
    for slide in slides_content:
        slide_index = slide['slide_index']
        content = slide['content']
        slide_folder = os.path.join(category_path, f"Slide_{slide_index + 1}")

        if not os.path.exists(slide_folder):
            os.makedirs(slide_folder)

        for item in content:
            if item['type'] == 'image':
                image_url = item['content']
                image_name = os.path.basename(image_url)
                image_path = os.path.join('static', image_url)
                dest_image_path = os.path.join(slide_folder, image_name)

                shutil.copy(image_path, dest_image_path)

        with open(os.path.join(slide_folder, f"Slide_{slide_index + 1}.txt"), 'w', encoding='utf-8') as f:
            f.write("\n".join([item['content'] for item in content if item['type'] == 'text']))


def extract_text_and_images_from_pptx(pptx_file, images_dir):
    prs = Presentation(pptx_file)
    slides_content = []

    for slide_index, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_content.append({'type': 'text', 'content': shape.text.strip()})
            if hasattr(shape, 'image'):
                image = shape.image
                image_file = os.path.join(images_dir, f'slide_{slide_index + 1}_{shape.shape_id}.png')
                image_bytes = image.blob
                with open(image_file, 'wb') as f:
                    f.write(image_bytes)
                slide_content.append({'type': 'image', 'content': image_file})
        slides_content.append({'slide_index': slide_index, 'content': slide_content})

    return slides_content


if __name__ == '__main__':
    create_directory(app.config['UPLOAD_FOLDER'])
    index_files(app.config['UPLOAD_FOLDER'])
    app.run(debug=True)
