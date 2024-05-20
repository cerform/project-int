import os
from werkzeug.utils import secure_filename
import mammoth
from pptx import Presentation
import fitz

class FileManager:
    def __init__(self, upload_folder):
        self.upload_folder = upload_folder

    def allowed_file(self, filename):
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'docx', 'pptx'}

    def get_categories(self):
        categories = {}
        for root, dirs, files in os.walk(self.upload_folder):
            for d in dirs:
                categories[d] = []
            for f in files:
                category = os.path.basename(root)
                if category in categories:
                    categories[category].append(f)
        return categories

    def save_file(self, file, category):
        category_path = os.path.join(self.upload_folder, secure_filename(category))
        if not os.path.exists(category_path):
            os.makedirs(category_path)
        file_path = os.path.join(category_path, secure_filename(file.filename))
        file.save(file_path)

    def extract_text_from_txt(self, path):
        with open(path, 'r', encoding='utf-8') as file:
            return file.read()

    def extract_text_from_pdf(self, path):
        doc = fitz.open(path)
        text = []
        for page in doc:
            text.append(page.get_text())
        return "\n".join(text)

    def extract_text_from_docx(self, path):
        with open(path, 'rb') as file:
            document = mammoth.extract_raw_text(file)
        return document.value

    def extract_text_from_pptx(self, path):
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
